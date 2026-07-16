"""
Adds slides to RAG.pptx using images from demos/pics/done,
with white background and images centered taking most of the slide space.
"""
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Config
PPTX_PATH = r'C:\Users\ahelmy\Grind\RAG.pptx'
PICS_DIR = r'C:\code\Personal\ai-projects\demos\pics'
OUTPUT_PATH = r'C:\Users\ahelmy\Grind\RAG2.pptx'  # Override the original file

# Get sorted images by filename (filenames contain timestamps)
images = sorted(
    glob.glob(os.path.join(PICS_DIR, '*.jpeg')) +
    glob.glob(os.path.join(PICS_DIR, '*.jpg')) +
    glob.glob(os.path.join(PICS_DIR, '*.png'))
)
print(f"Found {len(images)} images")

prs = Presentation(PPTX_PATH)
slide_w = prs.slide_width   # 13.33" = 12192000 EMU
slide_h = prs.slide_height  # 7.50"  = 6858000  EMU

# Use blank layout
blank_layout = prs.slide_layouts[6]  # Blank layout

# Full slide margins (leave small padding)
SLIDE_MARGIN = Inches(0.3)
CONTENT_LEFT   = SLIDE_MARGIN
CONTENT_TOP    = SLIDE_MARGIN
CONTENT_WIDTH  = slide_w - (2 * SLIDE_MARGIN)
CONTENT_HEIGHT = slide_h - (2 * SLIDE_MARGIN)

for i, img_path in enumerate(images):
    filename = os.path.basename(img_path)

    slide = prs.slides.add_slide(blank_layout)
    
    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    with Image.open(img_path) as im:
        img_w_px, img_h_px = im.size

    # Calculate scaled dimensions maintaining aspect ratio
    max_w_emu = CONTENT_WIDTH
    max_h_emu = CONTENT_HEIGHT
    img_ratio = img_w_px / img_h_px
    box_ratio = max_w_emu / max_h_emu

    if img_ratio > box_ratio:
        # Image is wider than box
        scaled_w = max_w_emu
        scaled_h = Emu(int(max_w_emu * img_h_px / img_w_px))
    else:
        # Image is taller than box
        scaled_h = max_h_emu
        scaled_w = Emu(int(max_h_emu * img_w_px / img_h_px))

    left_offset = CONTENT_LEFT + (max_w_emu - scaled_w) // 2
    top_offset  = CONTENT_TOP  + (max_h_emu - scaled_h) // 2

    slide.shapes.add_picture(img_path, left_offset, top_offset, scaled_w, scaled_h)

    print(f"  [{i+1:02d}/{len(images)}] Added: {filename}")

prs.save(OUTPUT_PATH)
print(f"\nSaved to: {OUTPUT_PATH}")
print(f"Total slides: {len(list(prs.slides))}")
