import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

# Config
PPTX_PATH = r'C:\Users\ahelmy\Grind\RAG.pptx'
PICS_DIR = r'C:\code\Personal\ai-projects\demos\pics'
OUTPUT_PATH = r'C:\Users\ahelmy\Grind\RAG_updated.pptx'

# Get sorted images
images = sorted(glob.glob(os.path.join(PICS_DIR, '*.jpeg')) + glob.glob(os.path.join(PICS_DIR, '*.jpg')) + glob.glob(os.path.join(PICS_DIR, '*.png')))
print(f"Found {len(images)} images")

prs = Presentation(PPTX_PATH)
slide_w = prs.slide_width   # 13.33 inches = 12192000 EMU
slide_h = prs.slide_height  # 7.50 inches  = 6858000 EMU

print(f"Slide size: {slide_w.inches:.2f} x {slide_h.inches:.2f} inches")

# Find the "Side bar Midnight/Dell Blue" layout (index 5)
sidebar_layout = None
for layout in prs.slide_master.slide_layouts:
    if layout.name == 'Side bar Midnight/Dell Blue':
        sidebar_layout = layout
        break

print(f"Using layout: {sidebar_layout.name if sidebar_layout else 'NOT FOUND'}")
print(f"Layout placeholders:")
for ph in sidebar_layout.placeholders:
    print(f"  idx={ph.placeholder_format.idx}, name={ph.name}, left={ph.left/914400:.2f}\", top={ph.top/914400:.2f}\", w={ph.width/914400:.2f}\", h={ph.height/914400:.2f}\"")

print(f"\nLayout shapes:")
for shape in sidebar_layout.shapes:
    print(f"  {shape.name}: type={shape.shape_type}, left={shape.left/914400:.2f}\", top={shape.top/914400:.2f}\", w={shape.width/914400:.2f}\", h={shape.height/914400:.2f}\"")
